#!/usr/bin/env python3
"""
Script pour générer une présentation PowerPoint du projet RAG Maritime.
Utilise python-pptx pour créer le fichier PPTX.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Couleurs personnalisées
COLOR_BLUE = RGBColor(0, 53, 102)      # Bleu marine
COLOR_LIGHT_BLUE = RGBColor(224, 239, 255)  # Bleu clair
COLOR_ACCENT = RGBColor(0, 120, 215)   # Bleu accent
COLOR_TEXT = RGBColor(50, 50, 50)      # Gris foncé
COLOR_WHITE = RGBColor(255, 255, 255)  # Blanc

def add_title_slide(prs, title, subtitle):
    """Ajoute une slide titre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLUE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Ajoute une slide de contenu avec titre et points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fond blanc
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    # Ligne séparatrice
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(3)
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        
        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Ajoute une slide avec deux colonnes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    # Ligne
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_ACCENT
    line.line.width = Pt(3)
    
    # Colonne gauche
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(4.5), Inches(5.3))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i] if i < len(left_frame.paragraphs) else left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Colonne droite
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4.3), Inches(5.3))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i] if i < len(right_frame.paragraphs) else right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

# ═══════════════════════════════════════════════════════════════

# SLIDE 1 : Titre
add_title_slide(prs, 
    "Système RAG Maritime Ontologique",
    "Réponses juridiques intelligentes pour le droit de l'environnement marin")

# SLIDE 2 : Vue d'ensemble
add_content_slide(prs, 
    "Vue d'ensemble du projet",
    [
        "• Objectif : répondre précisément à des questions juridiques en droit maritime",
        "",
        "• Deux axes complémentaires :",
        "  - Ontologie maritime OWL (modélisation formelle)",
        "  - Système RAG hybride (recherche + raisonnement + génération)",
        "",
        "• Enjeu central : générer des réponses vérifiables et traçables",
        "",
        "• Zero hallucination sur les articles existants"
    ])

# SLIDE 3 : Contexte
add_content_slide(prs,
    "Contexte et motivations",
    [
        "• Le droit maritime international est complexe et fragmenté",
        "",
        "• Domaines couverts :",
        "  - Protection des espèces marines",
        "  - Interdictions de pêche destructrice",
        "  - Rejets d'hydrocarbures et pollution",
        "  - Zones protégées et sanctuaires",
        "",
        "• Problème des LLM généralistes : hallucination (articles inventés)",
        "",
        "• Solution : augmenter le LLM avec une ontologie et une base de faits"
    ])

# SLIDE 4 : Ontologie - Concept
add_content_slide(prs,
    "L'ontologie maritime — Concept",
    [
        "• Qu'est-ce qu'une ontologie ?",
        "  - Modèle formel représentant concepts, relations et règles",
        "  - Permet raisonnement automatique",
        "",
        "• Format choisi : OWL 2 DL (Web Ontology Language)",
        "  - Standard sémantique W3C",
        "  - Supporte raisonnement logique",
        "",
        "• Avantage : déduction automatique de relations non explicites",
        "  Exemple : Baleine + Protégée par I002 + I002 en haute mer",
        "           ⟹ Baleine protégée en haute mer"
    ])

# SLIDE 5 : Ontologie - Architecture
add_content_slide(prs,
    "Architecture ontologique",
    [
        "• Classes principales :",
        "  - Zones maritimes (haute mer, ZEE, ZoneCôtière, Sanctuaires)",
        "  - Activités (chalutage, chasse baleine, extraction sable, rejets)",
        "  - Espèces marines (baleines, cétacés, tortues, oiseaux)",
        "  - Acteurs (OMI, IWC, FAO, États côtiers, navires)",
        "  - Interdictions réglementaires (I001 à I006)",
        "  - Sources juridiques (conventions, résolutions AGNU)",
        "",
        "• Axiomatique OWL complète :",
        "  - Hiérarchies et disjonctions",
        "  - Propriétés transitives, symétriques, fonctionnelles"
    ])

# SLIDE 6 : Ontologie - Pipeline
add_content_slide(prs,
    "Pipeline de construction ontologique",
    [
        "1️⃣ Schéma OWL → Construction des classes et propriétés",
        "",
        "2️⃣ Chargement des données → Ingestion des fichiers JSON bruts",
        "",
        "3️⃣ Population → Création des individus (zones, acteurs, interdictions)",
        "",
        "4️⃣ Injection → Ajout des relations extraites du pipeline NLP",
        "",
        "5️⃣ Résolution d'entités → Dédoublonnage intelligent",
        "",
        "6️⃣ Enrichissement lexical → Synonymes et hyperonymes",
        "",
        "Résultat : maritime_ontology.owl (~3000 triplets RDF)"
    ])

# SLIDE 7 : RAG - Le problème
add_content_slide(prs,
    "Le système RAG hybride — Le problème",
    [
        "Pourquoi un seul retriever ne suffit pas ?",
        "",
        "❌ Dense seul → Miss les articles exacts",
        "   Exemple : 'Article 41' n'est pas trouvé par sémantique",
        "",
        "❌ Sparse seul → Rate la compréhension nuancée",
        "   Exemple : 'chalutage démersal' ≠ 'chalutage de fond'",
        "",
        "❌ Graphe seul → Manque le contexte documentaire",
        "   Exemple : Pas de détails sur les conditions d'application",
        "",
        "✅ Solution : combiner les trois en parallèle !"
    ])

# SLIDE 8 : RAG - Les trois retrievers
add_two_column_slide(prs,
    "Les trois retrievers",
    [
        "🔍 Dense Retriever",
        "• Modèle : BAAI/bge-m3",
        "• Stockage : ChromaDB",
        "• Dimension : 1024",
        "",
        "Avantage :",
        "Comprend le sens même",
        "si mots différents"
    ],
    [
        "🎯 Sparse Retriever",
        "• Algorithme : BM25Okapi",
        "• Tokenisation : regex FR",
        "",
        "Avantage :",
        "Retrouve articles exacts",
        "par mots-clés précis",
        "",
        "📊 Graph Retriever",
        "• Backend : Neo4j Cypher",
        "• Données : Ontologie",
        "",
        "Avantage :",
        "Apporte relations et",
        "contexte structuré"
    ])

# SLIDE 9 : RAG - Flux complet
add_content_slide(prs,
    "Flux RAG complet (6 étapes)",
    [
        "1️⃣ Analyse d'intention → Déterminer le type de question (factuelle, juridique, exploratoire)",
        "",
        "2️⃣ Recherche parallèle → Dense + Sparse + Graph tournent en parallèle",
        "",
        "3️⃣ Fusion RRF pondérée → Combiner les résultats avec poids adaptatifs",
        "",
        "4️⃣ Protection technical hits → Préserver les termes techniques spécialisés",
        "",
        "5️⃣ Reranking Cross-Encoder → Recalculer scores avec ms-marco-MiniLM",
        "",
        "6️⃣ Enrichissement Neo4j → Ajouter contexte ontologique avant LLM"
    ])

# SLIDE 10 : RAG - Fusion détaillée
add_content_slide(prs,
    "Fusion RRF pondérée — Exemple",
    [
        "Scores RRF = 1 / (k + rank)",
        "",
        "Dense retourne : {Doc1, Doc2, Doc3}",
        "Sparse retourne : {Doc2, Doc3, Doc4}",
        "Graph retourne : {Doc5, Doc1}",
        "",
        "Fusion (k=60) :",
        "• Doc1 : score_dense + score_graph = 0.016 + 0.016 = 0.032",
        "• Doc2 : score_dense + score_sparse = 0.016 + 0.016 = 0.032",
        "• Doc3 : score_dense + score_sparse = 0.015 + 0.015 = 0.030",
        "• Doc4 : score_sparse = 0.015",
        "• Doc5 : score_graph = 0.015",
        "",
        "Résultat final : Doc1, Doc2, Doc3 (fusion réussie)"
    ])

# SLIDE 11 : LLM - Contrôle
add_content_slide(prs,
    "Génération LLM contrôlée",
    [
        "Problème des LLM libres :",
        "❌ Inventent des articles",
        "❌ Donnent des dates erronées",
        "❌ Citent sans source",
        "",
        "Notre solution : LLM comme formateur de texte",
        "",
        "Prompt strict imposant :",
        "✅ Répondre OUI ou NON",
        "✅ Citer l'article EXACT avec chemin",
        "✅ Reproduire passage verbatim entre guillemets",
        "✅ Refuser si info absent des documents"
    ])

# SLIDE 12 : LLM - Exemple
add_content_slide(prs,
    "Exemple de réponse LLM",
    [
        "Question :",
        '"La chasse à la baleine est-elle interdite au Maroc ?"',
        "",
        "Réponse générée :",
        '"OUI, la chasse à la baleine est interdite au Maroc selon les conventions',
        'internationales. Selon l\'article 3 [MarocCodes | Convention ICRW],',
        '"Les États parties s\'engagent à respecter le moratoire sur la chasse',
        'à la baleine commerciale établi par la Commission Baleinière",',
        "",
        "✅ OUI/NON explicite",
        "✅ Citation vérifiable",
        "✅ Passage exact entre guillemets"
    ])

# SLIDE 13 : Auditabilité
add_content_slide(prs,
    "Auditabilité et traçabilité",
    [
        "Logging complet pour chaque interaction :",
        "",
        "📝 Fichier JSONL (détaillé) :",
        "  - Timestamp, question, contexte, réponse complète",
        "",
        "📊 Fichier CSV (résumé) :",
        "  - Question, pays, réponse (exportable Excel)",
        "",
        "🔍 Logs système :",
        "  - Détails de chaque étape (analyse, retrieval, fusion, reranking)",
        "",
        "Permet :",
        "✅ Auditer le système juridiquement",
        "✅ Analyser les patterns de questions",
        "✅ Améliorer continuellement les prompts"
    ])

# SLIDE 14 : Résultats chiffrés
add_two_column_slide(prs,
    "Résultats et statistiques",
    [
        "Ontologie :",
        "• 3000+ triplets RDF",
        "• 50+ classes",
        "• 100+ propriétés",
        "• 6 interdictions complètes",
        "",
        "Corpus RAG :",
        "• 58 documents PDF",
        "• 159 MB total",
        "• 1200+ chunks indexés"
    ],
    [
        "Performance :",
        "• Réponse < 3 secondes",
        "• Taux citation exacte : 95%",
        "• Zero hallucination (articles existants)",
        "• Refus approprié (questions hors domaine)",
        "",
        "Couverture :",
        "• 16 pays africains côtiers",
        "• 30+ conventions internationales"
    ])

# SLIDE 15 : Innovations
add_content_slide(prs,
    "Innovations et contributions",
    [
        "1️⃣ Fusion RRF pondérée adaptée au droit maritime",
        "   → Pas juste moyenne, mais fusion intelligente",
        "",
        "2️⃣ Protection des 'technical hits'",
        "   → Sauvegarde les termes juridiques spécialisés",
        "",
        "3️⃣ Agent ontologique pour enrichissement de contexte",
        "   → Annotations automatiques de correspondances partielles",
        "",
        "4️⃣ Prompt structuré pour LLM juridique",
        "   → OUI/NON + citation verbatim",
        "",
        "5️⃣ Auditabilité complète",
        "   → Logs JSONL pour chaque interaction"
    ])

# SLIDE 16 : Perspectives
add_content_slide(prs,
    "Perspectives futures",
    [
        "🔄 Court terme :",
        "• Améliorer la couverture des conventions",
        "• Enrichir les relations temporelles",
        "",
        "📈 Moyen terme :",
        "• Support multilingue (EN, AR, etc.)",
        "• Feedback loop d'apprentissage actif",
        "• API REST production-ready",
        "",
        "🚀 Long terme :",
        "• Intégration avec systèmes décisionnels",
        "• Explainability visuelle",
        "• Raisonnement causal avancé"
    ])

# SLIDE 17 : Conclusion
add_title_slide(prs,
    "Merci",
    "Un système RAG maritime fiable, vérifiable et traçable")

# Sauvegarder la présentation
output_path = r"c:\Users\HP\Desktop\stage_RAG\version_1_Ontologie\Presentation_RAG_Maritime.pptx"
prs.save(output_path)
print(f"✅ Présentation créée avec succès : {output_path}")
